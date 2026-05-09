#!/usr/bin/env python3
"""
pptx_portfolio_generator.py
Generates a Portfolio/Showcase PowerPoint presentation for CMP Management Asia.

Structure:
  1. Title slide    — company name + presentation title
  2. Project slides — one per selected project (name, hero image, description)
  3. Closing slide  — tagline + contact info
"""

import json
import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

BASE_DIR = Path(__file__).parent
DATA_FILE = BASE_DIR / "data" / "projects.json"

# Slide dimensions — widescreen 16:9
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# Palette  (Midnight Executive + warm accent)
CLR_DARK   = RGBColor(0x1E, 0x27, 0x61)   # navy   — dark bg
CLR_MID    = RGBColor(0x2B, 0x37, 0x80)   # slightly lighter navy
CLR_ACCENT = RGBColor(0xE8, 0xA8, 0x20)   # gold accent
CLR_LIGHT  = RGBColor(0xCA, 0xDC, 0xFC)   # ice blue text
CLR_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CLR_BODY   = RGBColor(0xE8, 0xEE, 0xFF)   # near-white for body text

# Margins
MARGIN = Inches(0.5)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _load_data() -> dict:
    if not DATA_FILE.is_file():
        raise FileNotFoundError(f"Cannot find {DATA_FILE}")
    with open(DATA_FILE, encoding="utf-8") as f:
        return json.load(f)


def _rgb_fill(shape, color: RGBColor):
    """Solid fill a shape with an RGB colour."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text, font_size, bold=False,
                 color=CLR_WHITE, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb


def _set_bg(slide, color: RGBColor):
    """Set solid slide background colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _find_image(rel_path: str) -> Path | None:
    """Resolve an image path relative to BASE_DIR; return None if missing."""
    full = BASE_DIR / rel_path
    return full if full.is_file() else None


def _add_accent_bar(slide, left, top, width, height=Inches(0.055)):
    """Thin horizontal gold accent bar."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        left, top, width, height
    )
    _rgb_fill(bar, CLR_ACCENT)
    bar.line.fill.background()  # no border


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def _build_title_slide(prs: Presentation, title: str, subtitle: str, company_name: str):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    _set_bg(slide, CLR_DARK)

    w, h = SLIDE_W, SLIDE_H

    # Left accent column
    col = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), h)
    _rgb_fill(col, CLR_ACCENT)
    col.line.fill.background()

    # Company name (top-left)
    _add_textbox(
        slide,
        left=Inches(0.45), top=Inches(0.35),
        width=Inches(8), height=Inches(0.55),
        text=company_name.upper(),
        font_size=13, bold=True, color=CLR_ACCENT,
    )

    # Horizontal rule below company name
    _add_accent_bar(slide, Inches(0.45), Inches(0.95), Inches(5.5))

    # Main title
    _add_textbox(
        slide,
        left=Inches(0.45), top=Inches(2.0),
        width=Inches(7.5), height=Inches(2.2),
        text=title,
        font_size=48, bold=True, color=CLR_WHITE,
    )

    # Subtitle
    _add_textbox(
        slide,
        left=Inches(0.45), top=Inches(4.3),
        width=Inches(7.5), height=Inches(0.7),
        text=subtitle,
        font_size=20, bold=False, color=CLR_LIGHT,
    )

    # Decorative bottom bar
    _add_accent_bar(slide, Inches(0), h - Inches(0.12), w)

    return slide


def _build_project_slide(prs: Presentation, project: dict, index: int, total: int):
    """Left: text panel. Right: hero image."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    _set_bg(slide, CLR_DARK)

    w, h = SLIDE_W, SLIDE_H

    # ---- Right: image panel ----
    img_panel_left = Inches(5.5)
    img_panel_w    = w - img_panel_left

    image_path = None
    for img_rel in project.get("images", []):
        candidate = _find_image(img_rel)
        if candidate:
            image_path = candidate
            break

    if image_path:
        try:
            pic = slide.shapes.add_picture(
                str(image_path),
                img_panel_left, Inches(0),
                img_panel_w, h,
            )
            # Dim overlay on image for readability
            overlay = slide.shapes.add_shape(
                1,
                img_panel_left, Inches(0), img_panel_w, h
            )
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = RGBColor(0x1E, 0x27, 0x61)
            overlay.fill.fore_color.theme_color  # trigger
            # Set transparency via XML (50%)
            sp_pr = overlay.fill._fill._element
            solid_fill = sp_pr.find(qn("a:solidFill"))
            if solid_fill is not None:
                srgb = solid_fill.find(qn("a:srgbClr"))
                if srgb is None:
                    srgb = etree.SubElement(solid_fill, qn("a:srgbClr"))
                    srgb.set("val", "1E2761")
                alpha = etree.SubElement(srgb, qn("a:alpha"))
                alpha.set("val", "50000")  # 50% opacity = 50000 out of 100000
            overlay.line.fill.background()
        except Exception:
            # If image fails, fill right panel with mid-navy
            panel = slide.shapes.add_shape(
                1, img_panel_left, Inches(0), img_panel_w, h
            )
            _rgb_fill(panel, CLR_MID)
            panel.line.fill.background()
    else:
        panel = slide.shapes.add_shape(
            1, img_panel_left, Inches(0), img_panel_w, h
        )
        _rgb_fill(panel, CLR_MID)
        panel.line.fill.background()

    # ---- Left: text panel ----
    text_panel_w = Inches(5.3)

    # Thin left accent
    col = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.12), h)
    _rgb_fill(col, CLR_ACCENT)
    col.line.fill.background()

    # Counter  "03 / 12"
    counter_text = f"{index:02d} / {total:02d}"
    _add_textbox(
        slide,
        left=Inches(0.28), top=Inches(0.3),
        width=Inches(2.5), height=Inches(0.45),
        text=counter_text,
        font_size=11, bold=False, color=CLR_ACCENT,
    )

    # Category tag
    category = project.get("category", "").upper()
    _add_textbox(
        slide,
        left=Inches(0.28), top=Inches(0.75),
        width=Inches(3.5), height=Inches(0.4),
        text=category,
        font_size=11, bold=True, color=CLR_ACCENT, italic=True,
    )

    # Gold accent bar
    _add_accent_bar(slide, Inches(0.28), Inches(1.2), Inches(4.8))

    # Project name
    name = project.get("name", "Untitled Project")
    _add_textbox(
        slide,
        left=Inches(0.28), top=Inches(1.35),
        width=Inches(4.9), height=Inches(1.6),
        text=name,
        font_size=26, bold=True, color=CLR_WHITE,
    )

    # Description (subtitle)
    description = project.get("description", "")
    _add_textbox(
        slide,
        left=Inches(0.28), top=Inches(3.0),
        width=Inches(4.9), height=Inches(0.6),
        text=description,
        font_size=14, bold=False, color=CLR_LIGHT,
    )

    # Thin separator
    sep = slide.shapes.add_shape(1, Inches(0.28), Inches(3.65), Inches(4.8), Inches(0.02))
    _rgb_fill(sep, RGBColor(0x5A, 0x6A, 0xA0))
    sep.line.fill.background()

    # Meta info: Client + Role
    client = project.get("client", "")
    role   = project.get("role", "")
    year   = project.get("year", "")

    if client:
        _add_textbox(
            slide,
            left=Inches(0.28), top=Inches(3.75),
            width=Inches(4.9), height=Inches(0.4),
            text=f"Client: {client}",
            font_size=11, bold=False, color=CLR_BODY,
        )
    if role:
        _add_textbox(
            slide,
            left=Inches(0.28), top=Inches(4.18),
            width=Inches(4.9), height=Inches(0.55),
            text=f"Role: {role}",
            font_size=11, bold=False, color=CLR_BODY,
        )
    if year:
        _add_textbox(
            slide,
            left=Inches(0.28), top=Inches(4.76),
            width=Inches(1.5), height=Inches(0.4),
            text=str(year),
            font_size=12, bold=True, color=CLR_ACCENT,
        )

    # Bottom accent bar
    _add_accent_bar(slide, Inches(0), h - Inches(0.12), w)

    return slide


def _build_closing_slide(prs: Presentation, company: dict):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    _set_bg(slide, CLR_DARK)

    w, h = SLIDE_W, SLIDE_H

    # Left accent column
    col = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), h)
    _rgb_fill(col, CLR_ACCENT)
    col.line.fill.background()

    # "THANK YOU" heading
    _add_textbox(
        slide,
        left=Inches(0.45), top=Inches(1.5),
        width=Inches(9), height=Inches(1.4),
        text="THANK YOU",
        font_size=54, bold=True, color=CLR_WHITE,
        align=PP_ALIGN.LEFT,
    )

    # Gold rule
    _add_accent_bar(slide, Inches(0.45), Inches(3.0), Inches(6))

    # Tagline
    tagline = company.get("tagline", "")
    _add_textbox(
        slide,
        left=Inches(0.45), top=Inches(3.2),
        width=Inches(9), height=Inches(0.8),
        text=tagline,
        font_size=18, bold=False, color=CLR_LIGHT, italic=True,
    )

    # Contact block
    company_name = company.get("name", "CMP Management Asia")
    contact_lines = [
        company_name,
        "www.cmpmanagementasia.com",
        "contact@cmpmanagementasia.com",
    ]
    contact_text = "\n".join(contact_lines)
    _add_textbox(
        slide,
        left=Inches(0.45), top=Inches(4.3),
        width=Inches(7), height=Inches(1.8),
        text=contact_text,
        font_size=13, bold=False, color=CLR_BODY,
    )

    # Bottom accent bar
    _add_accent_bar(slide, Inches(0), h - Inches(0.12), w)

    return slide


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def generate_portfolio_pptx(selection: dict, output_path: str) -> str:
    """
    Generate a Portfolio/Showcase PPTX presentation.

    Args:
        selection: dict with keys:
            presentation_type  (str)
            projects           (list of project IDs)
            team_members       (list of member names)
            title              (str)
            subtitle           (str)
        output_path: str path where the .pptx file will be saved

    Returns:
        Absolute path to the generated PPTX file.
    """
    data = _load_data()
    all_projects = {p["id"]: p for p in data.get("projects", [])}
    company      = data.get("company", {"name": "CMP Management Asia", "tagline": ""})

    title    = selection.get("title", "Project Portfolio")
    subtitle = selection.get("subtitle", company.get("name", ""))
    project_ids = selection.get("projects", [])

    # Resolve selected projects (preserve order, skip unknowns)
    selected_projects = [all_projects[pid] for pid in project_ids if pid in all_projects]
    total = len(selected_projects)

    # Build presentation
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. Title slide
    _build_title_slide(prs, title, subtitle, company.get("name", "CMP Management Asia"))

    # 2. Project slides
    for idx, project in enumerate(selected_projects, start=1):
        _build_project_slide(prs, project, idx, total)

    # 3. Closing slide
    _build_closing_slide(prs, company)

    # Save
    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))

    return str(out.resolve())


# ---------------------------------------------------------------------------
# CLI test runner
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    sample_selection = {
        "presentation_type": "portfolio",
        "projects": ["atlas", "quarter", "obk", "trp", "kbh", "boonNak"],
        "team_members": [],
        "title": "Project Portfolio 2026",
        "subtitle": "CMP Management Asia – Selected Works",
    }

    output = str(BASE_DIR / "output_portfolio.pptx")
    result = generate_portfolio_pptx(sample_selection, output)
    print(f"Generated: {result}")
