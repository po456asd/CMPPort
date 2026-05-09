#!/usr/bin/env python3
"""
pptx_annual_generator.py
Generates Annual/Status Report PPTX presentations for CMP Management Asia.

Structure:
  1. Title slide        – report title + year/date range
  2. Summary statistics – total projects, team size, date range
  3. Timeline slides    – projects grouped by year (1-2 per slide)
  4. Team snapshot      – selected team members with roles
  5. Closing slide      – next steps / contact info
"""

import json
import os
from collections import defaultdict
from datetime import datetime

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
DATA_FILE = os.path.join(BASE_DIR, "data", "projects.json")

# Slide dimensions (widescreen 16:9)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# Brand colours
C_DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
C_ACCENT = RGBColor(0xC9, 0xA9, 0x6E)   # gold
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xF4, 0xF4, 0xF4)
C_MID    = RGBColor(0x55, 0x55, 0x66)

FONT_FAMILY = "Calibri"


# ---------------------------------------------------------------------------
# Data helpers
# ---------------------------------------------------------------------------

def _load_data() -> dict:
    with open(DATA_FILE, encoding="utf-8") as f:
        return json.load(f)


def _resolve_image(rel_path: str) -> str | None:
    """Return absolute path to image if it exists, else None."""
    abs_path = os.path.join(BASE_DIR, rel_path)
    return abs_path if os.path.isfile(abs_path) else None


# ---------------------------------------------------------------------------
# Slide-building helpers
# ---------------------------------------------------------------------------

def _add_filled_slide(prs: Presentation, bg_color: RGBColor):
    """Add a blank slide and fill background with a solid colour."""
    blank_layout = prs.slide_layouts[6]  # layout index 6 = completely blank
    slide = prs.slides.add_slide(blank_layout)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    return slide


def _add_textbox(slide, left, top, width, height,
                 text, font_size, bold=False, color=None,
                 align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT_FAMILY
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txb


def _add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def _safe_add_image(slide, img_path, left, top, width, height):
    """Add image if path is valid, otherwise add a placeholder rect."""
    if img_path and os.path.isfile(img_path):
        try:
            slide.shapes.add_picture(img_path, left, top, width, height)
            return True
        except Exception:
            pass
    # Placeholder
    _add_rect(slide, left, top, width, height, RGBColor(0xCC, 0xCC, 0xCC))
    _add_textbox(slide, left, top, width, height,
                 "[ image ]", 10, color=C_MID, align=PP_ALIGN.CENTER)
    return False


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def _build_title_slide(prs: Presentation, title: str, subtitle: str, year_label: str):
    slide = _add_filled_slide(prs, C_DARK)

    # Accent bar left edge
    _add_rect(slide, Inches(0), Inches(0), Inches(0.35), SLIDE_H, C_ACCENT)

    # Company label (top-right)
    _add_textbox(slide,
                 Inches(0.6), Inches(0.3),
                 Inches(12), Inches(0.5),
                 "CMP MANAGEMENT ASIA",
                 10, bold=True, color=C_ACCENT,
                 align=PP_ALIGN.RIGHT)

    # Main title
    _add_textbox(slide,
                 Inches(0.8), Inches(2.2),
                 Inches(11.5), Inches(1.6),
                 title,
                 40, bold=True, color=C_WHITE,
                 align=PP_ALIGN.LEFT)

    # Subtitle / year
    _add_textbox(slide,
                 Inches(0.8), Inches(3.9),
                 Inches(11.5), Inches(0.7),
                 subtitle,
                 20, color=C_ACCENT,
                 align=PP_ALIGN.LEFT)

    # Year label
    _add_textbox(slide,
                 Inches(0.8), Inches(4.65),
                 Inches(6), Inches(0.5),
                 year_label,
                 14, color=C_LIGHT,
                 align=PP_ALIGN.LEFT)

    # Bottom rule
    _add_rect(slide, Inches(0.8), Inches(6.9), Inches(11.5), Inches(0.05), C_ACCENT)


def _build_summary_slide(prs: Presentation,
                         total_projects: int,
                         team_size: int,
                         year_range: str,
                         categories: dict):
    slide = _add_filled_slide(prs, C_LIGHT)

    # Header bar
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), C_DARK)
    _add_textbox(slide,
                 Inches(0.4), Inches(0.2),
                 Inches(12), Inches(0.7),
                 "SUMMARY STATISTICS",
                 22, bold=True, color=C_WHITE)

    # Three stat cards
    card_w = Inches(3.5)
    card_h = Inches(2.2)
    card_top = Inches(1.6)
    card_gap = Inches(0.45)

    stats = [
        ("Projects Completed", str(total_projects), "total selected"),
        ("Team Members",       str(team_size),       "active professionals"),
        ("Period Covered",     year_range,            "date range"),
    ]

    for i, (label, value, sub) in enumerate(stats):
        left = Inches(0.45) + i * (card_w + card_gap)
        _add_rect(slide, left, card_top, card_w, card_h, C_DARK)
        _add_textbox(slide, left + Inches(0.2), card_top + Inches(0.15),
                     card_w - Inches(0.4), Inches(0.4),
                     label.upper(), 9, bold=True, color=C_ACCENT)
        _add_textbox(slide, left + Inches(0.2), card_top + Inches(0.55),
                     card_w - Inches(0.4), Inches(0.9),
                     value, 36, bold=True, color=C_WHITE)
        _add_textbox(slide, left + Inches(0.2), card_top + Inches(1.6),
                     card_w - Inches(0.4), Inches(0.4),
                     sub, 10, color=C_LIGHT)

    # Category breakdown
    _add_textbox(slide,
                 Inches(0.45), Inches(4.1),
                 Inches(12), Inches(0.5),
                 "Projects by Category",
                 14, bold=True, color=C_DARK)

    # Vertical stacked bars
    cat_top = Inches(4.7)
    for cat, count in sorted(categories.items()):
        bar_w = Inches(min(count * 1.2, 11.0))
        _add_rect(slide, Inches(0.45), cat_top, bar_w, Inches(0.38), C_ACCENT)
        _add_textbox(slide,
                     bar_w + Inches(0.55), cat_top,
                     Inches(4), Inches(0.38),
                     f"{cat.capitalize()}  ({count})", 10, color=C_DARK)
        cat_top += Inches(0.5)


def _build_timeline_slide(prs: Presentation, year: int, projects_slice: list,
                           all_projects_map: dict):
    """One slide for 1 or 2 projects from a given year."""
    slide = _add_filled_slide(prs, C_WHITE)

    # Header bar
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), C_DARK)
    _add_textbox(slide,
                 Inches(0.4), Inches(0.15),
                 Inches(9), Inches(0.7),
                 f"TIMELINE  —  {year}",
                 22, bold=True, color=C_WHITE)
    # Year badge
    _add_rect(slide, Inches(11.0), Inches(0.15), Inches(1.9), Inches(0.7), C_ACCENT)
    _add_textbox(slide,
                 Inches(11.0), Inches(0.15),
                 Inches(1.9), Inches(0.7),
                 str(year), 24, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

    n = len(projects_slice)
    col_w = Inches(5.8) if n == 2 else Inches(10.0)
    col_gap = Inches(0.6)
    img_h  = Inches(3.0)
    start_left = Inches(0.5) if n == 2 else Inches(1.65)

    for i, proj_id in enumerate(projects_slice):
        proj = all_projects_map.get(proj_id)
        if not proj:
            continue

        left = start_left + i * (col_w + col_gap)
        top_img = Inches(1.2)

        # Image
        img_path = None
        for img_rel in proj.get("images", []):
            candidate = _resolve_image(img_rel)
            if candidate:
                img_path = candidate
                break

        _safe_add_image(slide, img_path, left, top_img, col_w, img_h)

        # Project name
        _add_textbox(slide,
                     left, top_img + img_h + Inches(0.1),
                     col_w, Inches(0.5),
                     proj.get("name", proj_id),
                     14, bold=True, color=C_DARK)

        # Description / role
        desc = proj.get("description", "")
        role = proj.get("role", "")
        detail = f"{desc}  |  {role}" if desc and role else desc or role
        _add_textbox(slide,
                     left, top_img + img_h + Inches(0.65),
                     col_w, Inches(0.5),
                     detail, 10, color=C_MID)

        # Client
        client = proj.get("client", "")
        if client:
            _add_textbox(slide,
                         left, top_img + img_h + Inches(1.1),
                         col_w, Inches(0.35),
                         f"Client: {client}", 9, color=C_MID)


def _build_team_slide(prs: Presentation, team_members: list, all_team_map: dict):
    slide = _add_filled_slide(prs, C_DARK)

    # Header
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), C_ACCENT)
    _add_textbox(slide,
                 Inches(0.4), Inches(0.15),
                 Inches(12), Inches(0.7),
                 "TEAM SNAPSHOT",
                 24, bold=True, color=C_DARK)

    if not team_members:
        _add_textbox(slide, Inches(0.5), Inches(2.0), Inches(12), Inches(1.0),
                     "No team members selected.", 16, color=C_LIGHT)
        return

    # Up to 4 members per slide, in a 2x2 or 1x4 grid
    card_w = Inches(2.8)
    card_h = Inches(3.8)
    cols = min(len(team_members), 4)
    total_w = cols * card_w + (cols - 1) * Inches(0.3)
    start_left = (SLIDE_W - total_w) / 2

    for i, name in enumerate(team_members[:4]):
        member = all_team_map.get(name)
        if not member:
            continue

        left = start_left + i * (card_w + Inches(0.3))
        top  = Inches(1.3)

        # Card background
        _add_rect(slide, left, top, card_w, card_h,
                  RGBColor(0x25, 0x25, 0x3E))

        # Accent top bar
        _add_rect(slide, left, top, card_w, Inches(0.08), C_ACCENT)

        # Name
        _add_textbox(slide,
                     left + Inches(0.15), top + Inches(0.2),
                     card_w - Inches(0.3), Inches(0.55),
                     member.get("name", name),
                     13, bold=True, color=C_WHITE)

        # Role
        _add_textbox(slide,
                     left + Inches(0.15), top + Inches(0.78),
                     card_w - Inches(0.3), Inches(0.4),
                     member.get("role", ""),
                     10, color=C_ACCENT)

        # Bio (truncated)
        bio = member.get("bio", "")
        if len(bio) > 160:
            bio = bio[:157] + "..."
        _add_textbox(slide,
                     left + Inches(0.15), top + Inches(1.25),
                     card_w - Inches(0.3), Inches(2.3),
                     bio, 9, color=C_LIGHT, wrap=True)


def _build_closing_slide(prs: Presentation, company: dict):
    slide = _add_filled_slide(prs, C_DARK)

    # Gold bar
    _add_rect(slide, Inches(0), Inches(0), Inches(0.4), SLIDE_H, C_ACCENT)

    # Company name
    _add_textbox(slide,
                 Inches(0.8), Inches(1.8),
                 Inches(11.5), Inches(1.0),
                 company.get("name", "CMP Management Asia").upper(),
                 28, bold=True, color=C_WHITE,
                 align=PP_ALIGN.LEFT)

    # Tagline
    tagline = company.get("tagline", "From Vision to Guaranteed Reality.")
    _add_textbox(slide,
                 Inches(0.8), Inches(2.9),
                 Inches(10.0), Inches(0.7),
                 tagline,
                 16, color=C_ACCENT,
                 align=PP_ALIGN.LEFT)

    # Next steps heading
    _add_textbox(slide,
                 Inches(0.8), Inches(3.9),
                 Inches(5), Inches(0.5),
                 "Next Steps",
                 14, bold=True, color=C_WHITE)

    next_steps = [
        "Schedule strategy review with stakeholders",
        "Finalise project pipeline for next period",
        "Distribute report to partners & clients",
    ]
    for j, step in enumerate(next_steps):
        _add_textbox(slide,
                     Inches(0.8), Inches(4.5) + Inches(j * 0.45),
                     Inches(11.0), Inches(0.4),
                     f"•  {step}",
                     11, color=C_LIGHT)

    # Contact line
    _add_textbox(slide,
                 Inches(0.8), Inches(6.6),
                 Inches(11.5), Inches(0.4),
                 "www.cmpmanagementasia.com",
                 10, color=C_ACCENT,
                 align=PP_ALIGN.LEFT)

    # Bottom rule
    _add_rect(slide, Inches(0.8), Inches(6.85), Inches(11.5), Inches(0.04), C_ACCENT)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def generate_annual_pptx(selection: dict, output_path: str) -> str:
    """
    Generate an Annual/Status Report PPTX.

    Args:
        selection: dict with keys:
            presentation_type  – should be "annual"
            projects           – list of project IDs to include
            team_members       – list of team member names to include
            title              – presentation title
            subtitle           – subtitle / company name
        output_path: path where the .pptx file should be saved

    Returns:
        Absolute path to the generated PPTX file.
    """
    data = _load_data()
    all_projects   = data.get("projects", [])
    all_team       = data.get("team", [])
    company        = data.get("company", {})

    # Build lookup maps
    projects_map = {p["id"]: p for p in all_projects}
    team_map     = {m["name"]: m for m in all_team}

    # Resolve selection lists
    selected_project_ids  = selection.get("projects", [])
    selected_member_names = selection.get("team_members", [])
    title    = selection.get("title", "Annual Report")
    subtitle = selection.get("subtitle", company.get("name", "CMP Management Asia"))

    # Determine year label
    current_year = str(datetime.now().year)
    year_label   = subtitle if subtitle and any(c.isdigit() for c in subtitle) else current_year

    # Filter projects to those selected
    selected_projects = [projects_map[pid] for pid in selected_project_ids
                         if pid in projects_map]

    # Year range string
    years = sorted({p["year"] for p in selected_projects}) if selected_projects else [int(current_year)]
    if len(years) == 1:
        year_range = str(years[0])
    else:
        year_range = f"{years[0]} – {years[-1]}"

    # Category counts
    categories = defaultdict(int)
    for p in selected_projects:
        categories[p.get("category", "other")] += 1

    # Group project IDs by year
    by_year = defaultdict(list)
    for p in selected_projects:
        by_year[p["year"]].append(p["id"])

    # --------------- Build presentation ---------------
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. Title slide
    _build_title_slide(prs, title, subtitle, year_label)

    # 2. Summary statistics
    _build_summary_slide(
        prs,
        total_projects=len(selected_projects),
        team_size=len(selected_member_names) if selected_member_names
                  else len(all_team),
        year_range=year_range,
        categories=dict(categories),
    )

    # 3. Timeline slides (1-2 projects per slide, grouped by year)
    for year in sorted(by_year.keys()):
        proj_ids = by_year[year]
        # Split into chunks of 2
        for chunk_start in range(0, len(proj_ids), 2):
            chunk = proj_ids[chunk_start:chunk_start + 2]
            _build_timeline_slide(prs, year, chunk, projects_map)

    # 4. Team snapshot
    members_for_slide = selected_member_names if selected_member_names \
                        else [m["name"] for m in all_team]
    _build_team_slide(prs, members_for_slide, team_map)

    # 5. Closing slide
    _build_closing_slide(prs, company)

    # Save
    output_path = os.path.abspath(output_path)
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"[annual] Saved: {output_path}")
    return output_path


# ---------------------------------------------------------------------------
# CLI smoke-test
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    import sys

    data = json.load(open(DATA_FILE, encoding="utf-8"))
    all_ids   = [p["id"] for p in data["projects"]]
    all_names = [m["name"] for m in data["team"]]

    sample_selection = {
        "presentation_type": "annual",
        "projects":          all_ids,
        "team_members":      all_names,
        "title":             "Annual Report 2025",
        "subtitle":          "CMP Management Asia",
    }

    out = os.path.join(BASE_DIR, "output", "annual_report_test.pptx")
    result = generate_annual_pptx(sample_selection, out)
    print(f"Generated: {result}")
    if sys.platform == "win32":
        os.startfile(result)
