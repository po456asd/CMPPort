#!/usr/bin/env python3
"""
pptx_team_generator.py
Generates a Team Profile PowerPoint presentation for CMP Management Asia.

Slide structure (US-004):
  1. Title slide           – company name + presentation title
  2. Company overview      – tagline, team size, company info
  3. One slide per member  – name, role, bio
  4. Featured projects     – 4-6 key projects with images (2-per-row grid)
  5. Closing slide         – contact info
"""

import json
import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Constants – colour palette (Midnight Executive)
# ---------------------------------------------------------------------------
C_DARK_BG   = RGBColor(0x1E, 0x27, 0x61)   # navy  – title / closing bg
C_LIGHT_BG  = RGBColor(0xF4, 0xF6, 0xFB)   # off-white – content slides
C_ACCENT    = RGBColor(0x1E, 0x27, 0x61)   # navy accent bar
C_TEXT_DARK = RGBColor(0x1E, 0x1E, 0x2E)   # near-black body text
C_TEXT_LITE = RGBColor(0xFF, 0xFF, 0xFF)   # white for dark backgrounds
C_MUTED     = RGBColor(0x64, 0x74, 0x8B)   # slate for captions / labels
C_GOLD      = RGBColor(0xCA, 0xDC, 0xFC)   # ice-blue accent

BASE_DIR  = Path(__file__).parent
DATA_FILE = BASE_DIR / "data" / "projects.json"

# Slide dimensions for LAYOUT_16x9 (in inches)
SLIDE_W = 13.33
SLIDE_H = 7.5

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _load_data() -> dict:
    with open(DATA_FILE, encoding="utf-8") as f:
        return json.load(f)


def _rgb_fill(shape, color: RGBColor) -> None:
    """Set solid fill on a shape."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, x, y, w, h, color: RGBColor):
    """Add a filled rectangle (no border)."""
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    _rgb_fill(shape, color)
    shape.line.fill.background()  # no border
    return shape


def _add_textbox(slide, x, y, w, h, text, font_size, bold=False,
                 color: RGBColor = None, align=PP_ALIGN.LEFT,
                 italic=False, wrap=True):
    """Add a text box and return the paragraph."""
    from pptx.util import Inches, Pt
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txb


def _add_image_safe(slide, img_path: str, x, y, w, h):
    """Add image if the file exists; draw a placeholder rectangle otherwise."""
    abs_path = BASE_DIR / img_path if not os.path.isabs(img_path) else Path(img_path)
    if abs_path.exists():
        try:
            slide.shapes.add_picture(
                str(abs_path),
                Inches(x), Inches(y), Inches(w), Inches(h)
            )
            return True
        except Exception:
            pass
    # Placeholder
    _add_rect(slide, x, y, w, h, RGBColor(0xCC, 0xCC, 0xCC))
    _add_textbox(slide, x + 0.1, y + h / 2 - 0.15, w - 0.2, 0.3,
                 "[image]", 8, color=RGBColor(0x80, 0x80, 0x80),
                 align=PP_ALIGN.CENTER)
    return False


def _slide_background(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def _add_title_slide(prs: Presentation, company: dict, title: str, subtitle: str):
    """Slide 1: Title slide – dark background, company + presentation title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _slide_background(slide, C_DARK_BG)

    # Gold accent bar at top
    _add_rect(slide, 0, 0, SLIDE_W, 0.07, C_GOLD)

    # Company name
    _add_textbox(
        slide, 0.7, 1.4, 8.6, 0.7,
        company.get("name", "CMP Management Asia"),
        font_size=14, bold=False, color=C_GOLD, align=PP_ALIGN.CENTER
    )

    # Presentation title – large
    _add_textbox(
        slide, 0.7, 2.0, 8.6, 1.1,
        title,
        font_size=40, bold=True, color=C_TEXT_LITE, align=PP_ALIGN.CENTER
    )

    # Subtitle
    if subtitle:
        _add_textbox(
            slide, 0.7, 3.15, 8.6, 0.55,
            subtitle,
            font_size=16, color=C_GOLD, align=PP_ALIGN.CENTER, italic=True
        )

    # Tagline
    tagline = company.get("tagline", "")
    if tagline:
        _add_textbox(
            slide, 0.7, 4.5, 8.6, 0.5,
            tagline,
            font_size=11, color=RGBColor(0xCA, 0xDC, 0xFC),
            align=PP_ALIGN.CENTER, italic=True
        )

    # Gold accent bar at bottom
    _add_rect(slide, 0, SLIDE_H - 0.07, SLIDE_W, 0.07, C_GOLD)


def _add_overview_slide(prs: Presentation, company: dict, team: list):
    """Slide 2: Company overview – tagline, team size, company info."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_background(slide, C_LIGHT_BG)

    # Left navy accent bar
    _add_rect(slide, 0, 0, 0.12, SLIDE_H, C_ACCENT)

    # Slide title
    _add_textbox(
        slide, 0.4, 0.22, 9.2, 0.65,
        "Company Overview",
        font_size=28, bold=True, color=C_ACCENT
    )

    # Divider line (thin rect)
    _add_rect(slide, 0.4, 0.9, 9.2, 0.04, C_GOLD)

    company_name = company.get("name", "CMP Management Asia")
    tagline      = company.get("tagline", "")

    # Company name
    _add_textbox(slide, 0.4, 1.05, 9.2, 0.5,
                 company_name, font_size=20, bold=True, color=C_TEXT_DARK)

    # Tagline
    _add_textbox(slide, 0.4, 1.55, 9.2, 0.55,
                 tagline, font_size=14, color=C_MUTED, italic=True)

    # Stats row – 3 callout boxes
    stats = [
        ("30+",  "Years of\nExperience"),
        (str(len(team)), "Expert\nTeam Members"),
        ("3",    "Service\nDisciplines"),
    ]
    box_w, box_h = 2.5, 1.5
    start_x = (SLIDE_W - len(stats) * box_w - (len(stats) - 1) * 0.35) / 2
    for i, (number, label) in enumerate(stats):
        bx = start_x + i * (box_w + 0.35)
        by = 2.45
        _add_rect(slide, bx, by, box_w, box_h, C_DARK_BG)
        _add_textbox(slide, bx, by + 0.1, box_w, 0.75,
                     number, font_size=36, bold=True, color=C_GOLD,
                     align=PP_ALIGN.CENTER)
        _add_textbox(slide, bx, by + 0.85, box_w, 0.55,
                     label, font_size=10, color=C_TEXT_LITE,
                     align=PP_ALIGN.CENTER)

    # Bottom note
    _add_textbox(slide, 0.4, 4.95, 9.2, 0.4,
                 "Serving clients across Asia with integrated construction and project management solutions.",
                 font_size=9, color=C_MUTED, align=PP_ALIGN.CENTER, italic=True)


def _add_team_member_slide(prs: Presentation, member: dict):
    """Slide per team member: name, role, bio."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_background(slide, C_LIGHT_BG)

    # Left navy accent bar
    _add_rect(slide, 0, 0, 0.12, SLIDE_H, C_ACCENT)

    # "TEAM MEMBER" label
    _add_textbox(
        slide, 0.4, 0.22, 4.0, 0.4,
        "TEAM MEMBER", font_size=9, bold=True, color=C_MUTED
    )

    # Member name
    _add_textbox(
        slide, 0.4, 0.62, 6.5, 0.85,
        member.get("name", ""),
        font_size=34, bold=True, color=C_ACCENT
    )

    # Role badge (filled rect + text)
    role_text = member.get("role", "")
    _add_rect(slide, 0.4, 1.52, 4.5, 0.4, C_ACCENT)
    _add_textbox(
        slide, 0.4, 1.52, 4.5, 0.4,
        role_text, font_size=13, bold=True, color=C_TEXT_LITE
    )

    # Divider
    _add_rect(slide, 0.4, 2.05, 6.5, 0.04, C_GOLD)

    # Bio text
    bio = member.get("bio", "")
    _add_textbox(
        slide, 0.4, 2.2, 6.5, 2.9,
        bio, font_size=14, color=C_TEXT_DARK, wrap=True
    )

    # Right decorative block (navy)
    _add_rect(slide, 7.5, 0.5, 2.3, SLIDE_H - 0.5, C_DARK_BG)
    # Initial / monogram
    initials = "".join(w[0].upper() for w in member.get("name", "?").split()[:2])
    _add_textbox(
        slide, 7.5, 1.8, 2.3, 1.5,
        initials, font_size=72, bold=True, color=C_GOLD,
        align=PP_ALIGN.CENTER
    )


def _add_featured_projects_slide(prs: Presentation, projects: list, base_dir: Path):
    """
    Featured projects carousel – up to 6 projects in a 3-column x 2-row grid.
    Each cell: project image + name.
    """
    # Take up to 6 projects; if fewer exist, use all of them
    featured = projects[:6] if len(projects) >= 6 else projects

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_background(slide, C_DARK_BG)

    # Top accent bar
    _add_rect(slide, 0, 0, SLIDE_W, 0.07, C_GOLD)

    # Title
    _add_textbox(
        slide, 0.4, 0.15, 9.2, 0.65,
        "Featured Projects",
        font_size=28, bold=True, color=C_TEXT_LITE
    )

    cols   = 3
    rows   = 2
    pad    = 0.18
    margin_x = 0.4
    margin_y = 1.0
    cell_w = (SLIDE_W - 2 * margin_x - (cols - 1) * pad) / cols
    cell_h = (SLIDE_H - margin_y - 0.6 - (rows - 1) * pad) / rows  # leave footer

    for idx, proj in enumerate(featured[:cols * rows]):
        col = idx % cols
        row = idx // cols
        cx  = margin_x + col * (cell_w + pad)
        cy  = margin_y + row * (cell_h + pad)

        img_h = cell_h - 0.38  # leave space for label
        # Image area
        first_img = proj.get("images", [None])[0]
        if first_img:
            _add_image_safe(slide, first_img, cx, cy, cell_w, img_h)
        else:
            _add_rect(slide, cx, cy, cell_w, img_h, RGBColor(0x40, 0x50, 0x80))

        # Project name label
        _add_rect(slide, cx, cy + img_h, cell_w, 0.35, RGBColor(0x28, 0x35, 0x70))
        _add_textbox(
            slide, cx + 0.05, cy + img_h + 0.02, cell_w - 0.1, 0.31,
            proj.get("name", ""),
            font_size=8, bold=True, color=C_TEXT_LITE, align=PP_ALIGN.LEFT
        )

    # Bottom accent bar
    _add_rect(slide, 0, SLIDE_H - 0.07, SLIDE_W, 0.07, C_GOLD)


def _add_closing_slide(prs: Presentation, company: dict):
    """Closing slide: contact info."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_background(slide, C_DARK_BG)

    # Top + bottom gold bars
    _add_rect(slide, 0, 0, SLIDE_W, 0.07, C_GOLD)
    _add_rect(slide, 0, SLIDE_H - 0.07, SLIDE_W, 0.07, C_GOLD)

    # Company name
    _add_textbox(
        slide, 0.7, 0.8, 8.6, 0.7,
        company.get("name", "CMP Management Asia"),
        font_size=14, bold=False, color=C_GOLD, align=PP_ALIGN.CENTER
    )

    # Thank-you text
    _add_textbox(
        slide, 0.7, 1.45, 8.6, 1.0,
        "Thank You",
        font_size=44, bold=True, color=C_TEXT_LITE, align=PP_ALIGN.CENTER
    )

    # Tagline
    tagline = company.get("tagline", "")
    if tagline:
        _add_textbox(
            slide, 0.7, 2.5, 8.6, 0.55,
            tagline,
            font_size=13, color=C_GOLD, align=PP_ALIGN.CENTER, italic=True
        )

    # Contact block
    contact_lines = [
        "www.cmpmanagementasia.com",
        "info@cmpmanagementasia.com",
    ]
    _add_textbox(
        slide, 0.7, 3.35, 8.6, 0.45,
        contact_lines[0],
        font_size=12, color=RGBColor(0xCA, 0xDC, 0xFC), align=PP_ALIGN.CENTER
    )
    _add_textbox(
        slide, 0.7, 3.8, 8.6, 0.45,
        contact_lines[1],
        font_size=12, color=RGBColor(0xCA, 0xDC, 0xFC), align=PP_ALIGN.CENTER
    )


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def generate_team_pptx(selection: dict, output_path: str) -> str:
    """
    Generate a Team Profile PPTX.

    Args:
        selection: dict with keys:
            presentation_type  – str (expected "team")
            projects           – list of project IDs to feature (optional)
            team_members       – list of team member names to include
            title              – presentation title string
            subtitle           – subtitle / client name string
        output_path: str path where the .pptx file will be saved.

    Returns:
        Absolute path to the generated PPTX file.
    """
    data    = _load_data()
    all_projects = data.get("projects", [])
    all_team     = data.get("team", [])
    company      = data.get("company", {})

    title    = selection.get("title", "Our Team")
    subtitle = selection.get("subtitle", company.get("name", ""))

    # Resolve selected team members (preserve order from selection)
    member_names = selection.get("team_members", [])
    team_map = {m["name"]: m for m in all_team}
    if member_names:
        members = [team_map[n] for n in member_names if n in team_map]
    else:
        members = all_team  # default: all members

    # Resolve featured projects (4-6)
    project_ids = selection.get("projects", [])
    proj_map    = {p["id"]: p for p in all_projects}
    if project_ids:
        featured = [proj_map[pid] for pid in project_ids if pid in proj_map]
    else:
        # Auto-select: prefer featured=True first, then fill to 6
        featured = [p for p in all_projects if p.get("featured")]
        if len(featured) < 4:
            for p in all_projects:
                if p not in featured:
                    featured.append(p)
                if len(featured) >= 6:
                    break
    featured = featured[:6]

    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # 1. Title slide
    _add_title_slide(prs, company, title, subtitle)

    # 2. Company overview
    _add_overview_slide(prs, company, all_team)

    # 3. One slide per team member
    for member in members:
        _add_team_member_slide(prs, member)

    # 4. Featured projects carousel
    _add_featured_projects_slide(prs, featured, BASE_DIR)

    # 5. Closing slide
    _add_closing_slide(prs, company)

    output_path = os.path.abspath(output_path)
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    return output_path


# ---------------------------------------------------------------------------
# CLI test
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    import sys

    sample_selection = {
        "presentation_type": "team",
        "projects": [],          # auto-select featured
        "team_members": [
            "Johnny Beh",
            "Philaiwan Thathisa",
            "Worapot Kitcharenchai",
            "Clarence Sim B.H",
        ],
        "title": "Our Team",
        "subtitle": "CMP Management Asia",
    }

    out = os.path.join(BASE_DIR, "output", "team_profile.pptx")
    result = generate_team_pptx(sample_selection, out)
    print(f"PPTX saved: {result}")
